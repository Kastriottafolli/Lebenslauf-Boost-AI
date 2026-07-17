"""Sapphire-Nightfall Projektpraesentation (5 Folien) — passend zur aktuellen Website.

Ausgaben:
- lebenslauf-boost-ai-praesentation.pptx
- lebenslauf-boost-ai-praesentation.pdf
- slides/slide-01.png ... slide-05.png
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "presentation"
SLIDES = OUT / "slides"
ASSETS = ROOT / "docs" / "video" / "assets"
W, H = 1920, 1080

# Sapphire Nightfall — wie frontend/css/tokens.css
INK = "#262b40"
INK2 = "#2c444c"
PAPER = "#eef3fa"
PAPER2 = "#e2eaf6"
SURFACE = "#ffffff"
ACCENT = "#0474c4"
ACCENT2 = "#06457f"
BRIGHT = "#5379ae"
SOFT = "#a8c4ec"
MUTED = "#6b7f9c"
LINE = "#d5e0ef"
GOOD = "#1f7a5c"
WARN = "#9a6b1a"


def font(size: int, bold: bool = False, serif: bool = False) -> ImageFont.FreeTypeFont:
    if serif:
        paths = [
            Path("/System/Library/Fonts/Supplemental/Georgia Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Georgia.ttf"),
        ]
    else:
        paths = [
            Path("/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf"),
        ]
    for path in paths:
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def gradient_bg() -> Image.Image:
    """Dunkler Hero-Gradient wie auf der Website."""
    image = Image.new("RGB", (W, H), INK)
    draw = ImageDraw.Draw(image)
    for y in range(H):
        t = y / H
        r = int(0x26 + (0x06 - 0x26) * t * 0.55 + (0x2c - 0x26) * t * 0.35)
        g = int(0x2b + (0x45 - 0x2b) * t * 0.55 + (0x44 - 0x2b) * t * 0.25)
        b = int(0x40 + (0x7f - 0x40) * t * 0.55 + (0x4c - 0x40) * t * 0.2)
        draw.line((0, y, W, y), fill=(r, g, b))
    # soft blue orbs
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.ellipse((-280, -200, 620, 700), fill=(4, 116, 196, 40))
    od.ellipse((1280, 420, 2100, 1280), fill=(168, 196, 236, 35))
    return Image.alpha_composite(image.convert("RGBA"), overlay).convert("RGB")


def light_bg() -> Image.Image:
    image = Image.new("RGB", (W, H), PAPER)
    draw = ImageDraw.Draw(image)
    for y in range(H):
        t = y / H
        r = int(0xee + (0xe2 - 0xee) * t)
        g = int(0xf3 + (0xea - 0xf3) * t)
        b = int(0xfa + (0xf6 - 0xfa) * t)
        draw.line((0, y, W, y), fill=(r, g, b))
    return image


def wrap(draw: ImageDraw.ImageDraw, text: str, selected: ImageFont.ImageFont, width: int) -> list[str]:
    lines: list[str] = []
    current = ""
    for word in text.split():
        candidate = f"{current} {word}".strip()
        if draw.textlength(candidate, font=selected) <= width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def text_block(draw, xy, text, *, size, color=INK, bold=False, serif=False, width=800, spacing=10) -> int:
    selected = font(size, bold, serif)
    x, y = xy
    for line in wrap(draw, text, selected, width):
        draw.text((x, y), line, font=selected, fill=color)
        y += size + spacing
    return y


def badge(draw, x, y, label, *, fill=ACCENT, text="#fff") -> int:
    selected = font(15, True)
    w = int(draw.textlength(label, font=selected)) + 28
    draw.rounded_rectangle((x, y, x + w, y + 34), radius=17, fill=fill)
    draw.text((x + w // 2, y + 17), label, anchor="mm", font=selected, fill=text)
    return w


def pill(draw, x, y, label, *, outline=ACCENT, fill=SURFACE, text=ACCENT2) -> int:
    selected = font(15, True)
    w = int(draw.textlength(label, font=selected)) + 30
    draw.rounded_rectangle((x, y, x + w, y + 36), radius=18, fill=fill, outline=outline, width=2)
    draw.text((x + w // 2, y + 18), label, anchor="mm", font=selected, fill=text)
    return w


def card(draw, box, *, fill=SURFACE, outline=LINE, radius=20) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=2)


def footer_dark(draw) -> None:
    draw.line((80, 1015, 1840, 1015), fill=(255, 255, 255, 40), width=1)
    draw.text((80, 1035), "LEBENSLAUF BOOST AI", font=font(13, True), fill=SOFT)
    draw.text((1840, 1035), "Sapphire Nightfall · FastAPI · SQLite · Claude + OpenAI", anchor="ra", font=font(13), fill="#9eb6d4")


def footer_light(draw) -> None:
    draw.line((80, 1015, 1840, 1015), fill=LINE, width=2)
    draw.text((80, 1035), "LEBENSLAUF BOOST AI", font=font(13, True), fill=ACCENT2)
    draw.text((1840, 1035), "Sapphire Nightfall · FastAPI · SQLite · Claude + OpenAI", anchor="ra", font=font(13), fill=MUTED)


def fit_shot(source: Path, target: Image.Image, box: tuple[int, int, int, int]) -> None:
    x1, y1, x2, y2 = box
    width, height = x2 - x1, y2 - y1
    with Image.open(source).convert("RGB") as image:
        ratio = max(width / image.width, height / image.height)
        image = image.resize((int(image.width * ratio), int(image.height * ratio)), Image.Resampling.LANCZOS)
        left = (image.width - width) // 2
        top = max(0, (image.height - height) // 8)  # bias toward top of UI
        image = image.crop((left, top, left + width, top + height))
        # soft shadow
        shadow = Image.new("RGBA", (width + 24, height + 24), (0, 0, 0, 0))
        sd = ImageDraw.Draw(shadow)
        sd.rounded_rectangle((8, 10, width + 8, height + 12), radius=16, fill=(38, 43, 64, 55))
        shadow = shadow.filter(ImageFilter.GaussianBlur(10))
        target.paste(shadow, (x1 - 8, y1 - 6), shadow)
        mask = Image.new("L", (width, height), 0)
        ImageDraw.Draw(mask).rounded_rectangle((0, 0, width, height), radius=14, fill=255)
        target.paste(image, (x1, y1), mask)


def slide_1() -> Image.Image:
    image = gradient_bg()
    draw = ImageDraw.Draw(image)
    badge(draw, 80, 56, "PROJEKTPRÄSENTATION", fill=ACCENT)
    draw.text((1840, 72), "01", anchor="ra", font=font(22, True), fill=SOFT)

    draw.text((90, 200), "Lebenslauf", font=font(78, True, serif=True), fill="#f0f5fc")
    draw.text((90, 295), "Boost AI", font=font(86, True, serif=True), fill=SOFT)
    text_block(
        draw, (96, 420),
        "Die aktuelle Web-App: Sapphire Nightfall Design, Boosti-Tour, RAG, Claude & OpenAI, ATS-Check und sechs Export-Designs — strikt auf echten CV-Daten.",
        size=26, color="#c9d8ec", width=740, spacing=12,
    )
    x = 96
    for label in ["Boosti", "RAG", "2 KI-Modelle", "6 Designs", "BYOK"]:
        x += pill(draw, x, 620, label, outline=SOFT, fill=(38, 43, 64), text=SOFT) + 12

    fit_shot(ASSETS / "01-hero.png", image, (920, 140, 1840, 920))
    footer_dark(draw)
    return image


def slide_2() -> Image.Image:
    image = light_bg()
    draw = ImageDraw.Draw(image)
    badge(draw, 80, 56, "PROBLEM & LÖSUNG", fill=ACCENT2)
    draw.text((1840, 72), "02", anchor="ra", font=font(22, True), fill=BRIGHT)
    draw.text((80, 120), "Warum dieses Projekt?", font=font(44, True, serif=True), fill=INK)
    draw.text((82, 180), "Bewerbungen sind individuell — Lebensläufe oft nicht.", font=font(22), fill=MUTED)

    card(draw, (80, 250, 900, 900), fill=SURFACE, outline="#f0c9c4")
    draw.rounded_rectangle((80, 250, 900, 320), radius=20, fill="#b03d32")
    draw.rectangle((80, 290, 900, 320), fill="#b03d32")
    draw.text((120, 285), "Das Problem", font=font(28, True), fill="#fff")
    bullets = [
        "Jede Stelle verlangt andere Schwerpunkte und Keywords.",
        "Manuelle Anpassung kostet Zeit und ist fehleranfällig.",
        "Generative KI kann Fähigkeiten oder Erfolge erfinden.",
        "ATS-Systeme filtern viele Lebensläufe vor dem Interview.",
    ]
    y = 360
    for item in bullets:
        draw.ellipse((120, y + 8, 134, y + 22), fill="#b03d32")
        y = text_block(draw, (150, y), item, size=22, color=INK2, width=700, spacing=8) + 18

    card(draw, (960, 250, 1840, 900), fill=SURFACE, outline="#b7dccb")
    draw.rounded_rectangle((960, 250, 1840, 320), radius=20, fill=GOOD)
    draw.rectangle((960, 290, 1840, 320), fill=GOOD)
    draw.text((1000, 285), "Die Lösung", font=font(28, True), fill="#fff")
    solutions = [
        "RAG liefert nur belegte CV-Abschnitte an die KI.",
        "Professionelle Prompts steuern Claude und OpenAI.",
        "5-Sekunden-Ladeübergänge machen Schritte sichtbar.",
        "ATS-Score, Verfeinern und 6 Designs bis zum Export.",
        "API-Key-Links zu Anthropic & OpenAI — oder Demo-Modus.",
    ]
    y = 360
    for item in solutions:
        draw.ellipse((1000, y + 8, 1014, y + 22), fill=GOOD)
        y = text_block(draw, (1030, y), item, size=22, color=INK2, width=760, spacing=8) + 16

    footer_light(draw)
    return image


def slide_3() -> Image.Image:
    image = light_bg()
    draw = ImageDraw.Draw(image)
    badge(draw, 80, 56, "WORKFLOW", fill=ACCENT)
    draw.text((1840, 72), "03", anchor="ra", font=font(22, True), fill=BRIGHT)
    draw.text((80, 120), "Vom Upload zum fertigen Lebenslauf", font=font(42, True, serif=True), fill=INK)
    draw.text((82, 178), "Drei UI-Schritte — wie auf der aktuellen Website.", font=font(22), fill=MUTED)

    shots = [
        (ASSETS / "02-input-upload.png", "1", "Eingabe", "Stelle · CV · Key-Links"),
        (ASSETS / "03-compare-output.png", "2", "Bearbeiten", "Vergleich · ATS · Refine"),
        (ASSETS / "05-design-export.png", "3", "Design", "6 Designs · PDF/Word"),
    ]
    xs = [80, 680, 1280]
    for (src, num, title, sub), x in zip(shots, xs):
        card(draw, (x, 240, x + 560, 900), fill=SURFACE, outline=LINE)
        draw.ellipse((x + 24, 262, x + 68, 306), fill=ACCENT)
        draw.text((x + 46, 284), num, anchor="mm", font=font(20, True), fill="#fff")
        draw.text((x + 84, 268), title, font=font(24, True), fill=INK)
        draw.text((x + 84, 300), sub, font=font(15), fill=MUTED)
        fit_shot(src, image, (x + 24, 340, x + 536, 860))
        if x != xs[-1]:
            draw.polygon([(x + 570, 560), (x + 595, 575), (x + 570, 590)], fill=ACCENT)

    footer_light(draw)
    return image


def slide_4() -> Image.Image:
    image = light_bg()
    draw = ImageDraw.Draw(image)
    badge(draw, 80, 56, "TECHNIK", fill=ACCENT2)
    draw.text((1840, 72), "04", anchor="ra", font=font(22, True), fill=BRIGHT)
    draw.text((80, 120), "Architektur der aktuellen App", font=font(42, True, serif=True), fill=INK)
    draw.text((82, 178), "Frontend, Backend, KI und Persistenz klar getrennt.", font=font(22), fill=MUTED)

    boxes = [
        (80, 260, 480, 620, "Frontend", ["Sapphire Nightfall", "Boosti-Tour DE/EN", "5s Lade-Overlay", "BYOK + Key-Links"], ACCENT),
        (520, 260, 920, 620, "Backend", ["FastAPI · Pydantic", "Session-Orchestrierung", "Validierung", "Export-Service"], ACCENT2),
        (960, 260, 1360, 620, "KI & RAG", ["Claude + OpenAI", "Professionelle Prompts", "Few-shot + CoT", "ATS-Keyword-Check"], BRIGHT),
        (1400, 260, 1840, 620, "Daten", ["SQLite + SQLAlchemy", "sessions", "cv_documents", "generations · messages"], INK2),
    ]
    for x1, y1, x2, y2, title, lines, color in boxes:
        card(draw, (x1, y1, x2, y2), fill=SURFACE, outline=LINE)
        draw.rounded_rectangle((x1, y1, x2, y1 + 70), radius=20, fill=color)
        draw.rectangle((x1, y1 + 40, x2, y1 + 70), fill=color)
        draw.text(((x1 + x2) // 2, y1 + 35), title, anchor="mm", font=font(24, True), fill="#fff")
        y = y1 + 110
        for line in lines:
            draw.ellipse((x1 + 32, y + 6, x1 + 44, y + 18), fill=ACCENT)
            draw.text((x1 + 58, y), line, font=font(20), fill=INK2)
            y += 48

    # pipeline strip
    card(draw, (80, 680, 1840, 900), fill=INK, outline=ACCENT2)
    draw.text((110, 720), "Request-Ablauf", font=font(22, True), fill=SOFT)
    flow = "Upload  →  RAG  →  Prompt  →  Claude/OpenAI  →  ATS  →  Refine  →  6 Designs  →  PDF/Word"
    draw.text((110, 790), flow, font=font(24, True), fill="#f0f5fc")
    draw.text((110, 850), "Keys bleiben im Browser. Ohne Key läuft der Demo-Modus.", font=font(18), fill="#9eb6d4")

    footer_light(draw)
    return image


def slide_5() -> Image.Image:
    image = gradient_bg()
    draw = ImageDraw.Draw(image)
    badge(draw, 80, 56, "ERGEBNIS", fill=ACCENT)
    draw.text((1840, 72), "05", anchor="ra", font=font(22, True), fill=SOFT)
    draw.text((80, 120), "Was die aktuelle Website liefert", font=font(42, True, serif=True), fill="#f0f5fc")
    draw.text((82, 178), "MVP mit professionellem UI und kontrollierter KI.", font=font(22), fill="#9eb6d4")

    cols = [
        ("Produkt", ACCENT, ["Sapphire Nightfall UI", "Boosti-geführte Tour", "Drei UI-Schritte", "6 Export-Designs", "DE / EN"]),
        ("KI-Qualität", BRIGHT, ["Professionelle Prompts", "RAG & Faktentreue", "Claude ↔ OpenAI", "ATS-Keyword-Check", "Conversation History"]),
        ("Nächste Schritte", ACCENT2, ["PostgreSQL + Alembic", "Login & Konten", "Versions-Rollback", "Anschreiben", "Tests & Docker"]),
    ]
    x = 80
    for title, color, items in cols:
        draw.rounded_rectangle((x, 250, x + 560, 780), radius=20, fill=(15, 28, 48), outline=SOFT, width=2)
        draw.rounded_rectangle((x, 250, x + 560, 330), radius=20, fill=color)
        draw.rectangle((x, 300, x + 560, 330), fill=color)
        draw.text((x + 280, 290), title, anchor="mm", font=font(26, True), fill="#fff")
        y = 380
        for item in items:
            draw.ellipse((x + 40, y + 8, x + 54, y + 22), fill=SOFT)
            draw.text((x + 72, y), item, font=font(22), fill="#e8f0fa")
            y += 58
        x += 600

    draw.rounded_rectangle((200, 820, 1720, 960), radius=24, fill=ACCENT, outline=SOFT, width=2)
    draw.text((960, 875), "Echte CV-Daten · kontrollierte KI · professioneller Export", anchor="mm", font=font(28, True), fill="#fff")
    draw.text((960, 925), "Relevant finden · transparent vergleichen · gezielt verbessern · herunterladen", anchor="mm", font=font(18), fill=SOFT)
    footer_dark(draw)
    return image


def save_pptx(slide_paths: list[Path], destination: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for path in slide_paths:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    presentation.save(destination)


def save_pdf(slide_paths: list[Path], destination: Path) -> None:
    page = landscape((540, 960))
    pdf = canvas.Canvas(str(destination), pagesize=page)
    for path in slide_paths:
        pdf.drawImage(str(path), 0, 0, width=960, height=540)
        pdf.showPage()
    pdf.save()


def main() -> None:
    SLIDES.mkdir(parents=True, exist_ok=True)
    renderers = [slide_1, slide_2, slide_3, slide_4, slide_5]
    paths: list[Path] = []
    for index, renderer in enumerate(renderers, 1):
        path = SLIDES / f"slide-{index:02d}.png"
        renderer().save(path, quality=95)
        paths.append(path)

    pptx = OUT / "lebenslauf-boost-ai-praesentation.pptx"
    pdf = OUT / "lebenslauf-boost-ai-praesentation.pdf"
    save_pptx(paths, pptx)
    save_pdf(paths, pdf)
    print(f"Erstellt: {pptx}")
    print(f"Erstellt: {pdf}")


if __name__ == "__main__":
    main()
